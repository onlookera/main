"""
DocMaster — PPT 转 Word 服务
提取 PPT 中所有文本内容，生成结构化的 Word 文档
纯 Python 方案，无需外部依赖
"""

import os
from pathlib import Path
from pptx import Presentation as PptxPresentation
from docx import Document
from docx.shared import Pt, Inches, RGBColor


def convert_ppt_to_word(
    input_path: str,
    output_path: str,
    progress_callback=None,
) -> str:
    """
    将 PPT 文本内容提取到 Word 文档

    处理策略:
        1. 遍历每张幻灯片
        2. 提取所有文本框中的文字
        3. 保留基本的层级结构和序号
        4. 输出干净的 Word 文档

    参数:
        input_path: 源 .pptx 文件路径
        output_path: 目标 .docx 文件路径
        progress_callback: 进度回调 callback(percent: int)

    返回:
        输出文件路径
    """
    if not os.path.exists(input_path):
        raise FileNotFoundError(f'输入文件不存在: {input_path}')

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)

    # 读取源 PPT
    ppt_prs = PptxPresentation(input_path)
    total_slides = len(ppt_prs.slides)

    if total_slides == 0:
        raise RuntimeError('PPT 文件没有幻灯片，无法转换')

    # 创建目标 Word 文档
    doc = Document()

    # 设置默认字体大小
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Microsoft YaHei'
    font.size = Pt(11)

    # 遍历所有幻灯片
    for slide_num, slide in enumerate(ppt_prs.slides):
        # 幻灯片标题
        heading = doc.add_heading(level=1)
        heading_run = heading.add_run(f'幻灯片 {slide_num + 1}')
        heading_run.font.size = Pt(16)
        heading_run.font.color.rgb = RGBColor(0x33, 0x41, 0x55)

        # 提取所有形状中的文本
        text_extracted = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # 根据缩进级别设置样式
                        level = para.level if para.level is not None else 0
                        if level == 0:
                            p = doc.add_paragraph()
                            run = p.add_run(text)
                            run.font.size = Pt(11)
                        else:
                            # 缩进级别用列表样式
                            p = doc.add_paragraph(style='List Bullet')
                            p.clear()
                            run = p.add_run(text)
                            run.font.size = Pt(10)

                        text_extracted = True

        # 如果幻灯片没有文本，标注
        if not text_extracted:
            p = doc.add_paragraph()
            run = p.add_run('[此幻灯片无可提取的文本内容]')
            run.font.size = Pt(10)
            run.font.italic = True
            run.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

        # 幻灯片间隔
        doc.add_paragraph('')

        if progress_callback:
            progress_callback(int((slide_num + 1) / total_slides * 100))

    # 保存 Word 文档
    doc.save(output_path)
    return output_path
