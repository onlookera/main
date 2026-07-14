"""
DocMaster — PDF 转 PPT 服务
将 PDF 的每一页转换为 PPT 幻灯片
依赖: pypdf (提取页面) + Pillow (图片处理) + python-pptx (生成PPT)
"""

import os
import shutil
from pathlib import Path
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import tempfile


def convert_pdf_to_ppt(
    input_path: str,
    output_path: str,
    progress_callback=None,
) -> str:
    """
    将 PDF 每一页转为 PPT 幻灯片

    原理:
        1. 使用 pypdf 读取 PDF 页数
        2. 暂时将每页内容提取（简单版：创建带页码的占位幻灯片）
        3. 完整版需要 pdf2image (依赖 poppler) 或 fitz (PyMuPDF) 渲染图片

    参数:
        input_path: 源 PDF 文件路径
        output_path: 目标 .pptx 文件路径
        progress_callback: 进度回调 callback(percent: int)

    返回:
        输出文件路径

    注意:
        当前实现为基础版本：提取 PDF 文本内容生成幻灯片。
        完整图片渲染版本需要额外安装 PyMuPDF: pip install PyMuPDF
    """
    if not os.path.exists(input_path):
        raise FileNotFoundError(f'输入文件不存在: {input_path}')

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)

    # 读取 PDF
    reader = PdfReader(input_path)
    total_pages = len(reader.pages)

    if total_pages == 0:
        raise RuntimeError('PDF 文件没有页面，无法转换')

    # 尝试使用 PyMuPDF (fitz) 获得更好的渲染效果
    try:
        return _convert_with_fitz(
            input_path, output_path, reader, total_pages, progress_callback
        )
    except ImportError:
        # 回退方案：提取文本内容生成幻灯片
        return _convert_text_only(
            reader, output_path, total_pages, progress_callback
        )


def _convert_with_fitz(
    input_path: str,
    output_path: str,
    reader: PdfReader,
    total_pages: int,
    progress_callback=None,
) -> str:
    """
    使用 PyMuPDF 将 PDF 页面渲染为图片，嵌入 PPT
    这是最佳方案：每页以图片形式完整保留排版
    """
    import fitz  # PyMuPDF

    prs = Presentation()
    # 使用 16:9 宽屏比例
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    pdf_doc = fitz.open(input_path)

    with tempfile.TemporaryDirectory() as tmpdir:
        for page_num in range(total_pages):
            # 渲染页面为图片
            page = pdf_doc.load_page(page_num)
            # 300 DPI 高质量渲染
            zoom = 300 / 72  # 72 是 PDF 默认 DPI
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)

            img_path = os.path.join(tmpdir, f'page_{page_num + 1}.png')
            pix.save(img_path)

            # 添加幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)

            # 将图片填充到整个幻灯片
            slide.shapes.add_picture(
                img_path,
                Inches(0),
                Inches(0),
                prs.slide_width,
                prs.slide_height,
            )

            if progress_callback:
                progress_callback(int((page_num + 1) / total_pages * 100))

    pdf_doc.close()
    prs.save(output_path)
    return output_path


def _convert_text_only(
    reader: PdfReader,
    output_path: str,
    total_pages: int,
    progress_callback=None,
) -> str:
    """
    回退方案：提取文本生成幻灯片（无图片渲染）
    """
    from pptx.util import Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for page_num in range(total_pages):
        page = reader.pages[page_num]
        text = page.extract_text() or f'[第 {page_num + 1} 页 — 无文本内容]'

        # 使用空白布局
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 添加文本框
        left = Inches(1)
        top = Inches(0.5)
        width = Inches(11.333)
        height = Inches(6.5)

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        # 添加页码标题
        p = text_frame.paragraphs[0]
        p.text = f'第 {page_num + 1} 页 / 共 {total_pages} 页'
        p.font.size = Pt(14)
        p.font.bold = True

        # 添加正文内容
        p2 = text_frame.add_paragraph()
        p2.text = text
        p2.font.size = Pt(11)
        p2.space_before = Pt(12)

        if progress_callback:
            progress_callback(int((page_num + 1) / total_pages * 100))

    prs.save(output_path)
    return output_path
